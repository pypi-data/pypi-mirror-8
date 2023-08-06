# encoding: utf-8

"""
Test suite for pptx.text module
"""

from __future__ import absolute_import, print_function

import pytest

from pptx.dml.color import ColorFormat
from pptx.dml.fill import FillFormat
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_UNDERLINE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.text import Font, _Hyperlink, _Paragraph, _Run, TextFrame
from pptx.util import Inches, Pt

from .oxml.unitdata.text import (
    a_bodyPr, a_txBody, a_p, a_t, an_hlinkClick, an_r, an_rPr
)
from .unitutil.cxml import element, xml
from .unitutil.mock import (
    class_mock, instance_mock, loose_mock, property_mock
)


class DescribeTextFrame(object):

    def it_knows_its_autosize_setting(self, autosize_get_fixture):
        text_frame, expected_value = autosize_get_fixture
        assert text_frame.auto_size == expected_value

    def it_can_change_its_autosize_setting(self, autosize_set_fixture):
        text_frame, value, expected_xml = autosize_set_fixture
        text_frame.auto_size = value
        assert text_frame._txBody.xml == expected_xml

    def it_knows_the_number_of_paragraphs_it_contains(
            self, txBody, txBody_with_2_paras):
        assert len(TextFrame(txBody, None).paragraphs) == 1
        assert len(TextFrame(txBody_with_2_paras, None).paragraphs) == 2

    def it_can_add_a_paragraph_to_the_text_it_contains(
            self, txBody, txBody_with_2_paras_xml):
        text_frame = TextFrame(txBody, None)
        text_frame.add_paragraph()
        assert text_frame._txBody.xml == txBody_with_2_paras_xml

    def it_knows_what_text_it_contains(self, text_get_fixture):
        text_frame, expected_value = text_get_fixture
        assert text_frame.text == expected_value

    def it_can_replace_the_text_it_contains(
            self, txBody, txBody_with_text_xml):
        text_frame = TextFrame(txBody, None)
        text_frame.text = 'foobar'
        assert txBody.xml == txBody_with_text_xml

    def it_knows_its_margin_settings(self, margin_get_fixture):
        text_frame, prop_name, unit, expected_value = margin_get_fixture
        margin_value = getattr(text_frame, prop_name)
        assert getattr(margin_value, unit) == expected_value

    def it_can_change_its_margin_settings(self, margin_set_fixture):
        text_frame, prop_name, new_value, expected_xml = margin_set_fixture
        setattr(text_frame, prop_name, new_value)
        assert text_frame._txBody.xml == expected_xml

    def it_raises_on_attempt_to_set_margin_to_non_int(self, text_frame):
        with pytest.raises(TypeError):
            text_frame.margin_bottom = '0.1'

    def it_can_change_its_vertical_anchor_setting(
            self, txBody, txBody_with_anchor_ctr_xml):
        text_frame = TextFrame(txBody, None)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        assert text_frame._txBody.xml == txBody_with_anchor_ctr_xml

    def it_can_change_the_word_wrap_setting(
            self, txBody, txBody_with_wrap_on_xml, txBody_with_wrap_off_xml,
            txBody_xml):
        text_frame = TextFrame(txBody, None)
        assert text_frame.word_wrap is None

        text_frame.word_wrap = True
        assert text_frame._txBody.xml == txBody_with_wrap_on_xml
        assert text_frame.word_wrap is True

        text_frame.word_wrap = False
        assert text_frame._txBody.xml == txBody_with_wrap_off_xml
        assert text_frame.word_wrap is False

        text_frame.word_wrap = None
        assert text_frame._txBody.xml == txBody_xml
        assert text_frame.word_wrap is None

    def it_knows_the_part_it_belongs_to(self, text_frame_with_parent_):
        text_frame, parent_ = text_frame_with_parent_
        part = text_frame.part
        assert part is parent_.part

    # fixtures ---------------------------------------------

    @pytest.fixture(params=[
        ('p:txBody/a:bodyPr', None),
        ('p:txBody/a:bodyPr/a:noAutofit',   MSO_AUTO_SIZE.NONE),
        ('p:txBody/a:bodyPr/a:spAutoFit',   MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT),
        ('p:txBody/a:bodyPr/a:normAutofit', MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE),
    ])
    def autosize_get_fixture(self, request):
        txBody_cxml, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        return text_frame, expected_value

    @pytest.fixture(params=[
        ('p:txBody/a:bodyPr',               MSO_AUTO_SIZE.NONE,
         'p:txBody/a:bodyPr/a:noAutofit'),
        ('p:txBody/a:bodyPr/a:noAutofit',   MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
         'p:txBody/a:bodyPr/a:spAutoFit'),
        ('p:txBody/a:bodyPr/a:spAutoFit',   MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
         'p:txBody/a:bodyPr/a:normAutofit'),
        ('p:txBody/a:bodyPr/a:normAutofit', None,
         'p:txBody/a:bodyPr'),
    ])
    def autosize_set_fixture(self, request):
        txBody_cxml, value, expected_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        expected_xml = xml(expected_cxml)
        return text_frame, value, expected_xml

    @pytest.fixture(params=[
        ('p:txBody/a:bodyPr',             'left',   'emu',    Inches(0.1)),
        ('p:txBody/a:bodyPr',             'top',    'emu',    Inches(0.05)),
        ('p:txBody/a:bodyPr',             'right',  'emu',    Inches(0.1)),
        ('p:txBody/a:bodyPr',             'bottom', 'emu',    Inches(0.05)),
        ('p:txBody/a:bodyPr{lIns=9144}',  'left',   'cm',     0.0254),
        ('p:txBody/a:bodyPr{tIns=18288}', 'top',    'mm',     0.508),
        ('p:txBody/a:bodyPr{rIns=76200}', 'right',  'pt',     6.0),
        ('p:txBody/a:bodyPr{bIns=36576}', 'bottom', 'inches', 0.04),
    ])
    def margin_get_fixture(self, request):
        txBody_cxml, side, unit, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        prop_name = "margin_%s" % side
        return text_frame, prop_name, unit, expected_value

    @pytest.fixture(params=[
        ('p:txBody/a:bodyPr',            'left',   Inches(0.11),
         'p:txBody/a:bodyPr{lIns=100584}'),
        ('p:txBody/a:bodyPr{tIns=1234}', 'top',    Inches(0.12),
         'p:txBody/a:bodyPr{tIns=109728}'),
        ('p:txBody/a:bodyPr{rIns=2345}', 'right',  Inches(0.13),
         'p:txBody/a:bodyPr{rIns=118872}'),
        ('p:txBody/a:bodyPr{bIns=3456}', 'bottom', Inches(0.14),
         'p:txBody/a:bodyPr{bIns=128016}'),
        ('p:txBody/a:bodyPr', 'left',   Inches(0.1),  'p:txBody/a:bodyPr'),
        ('p:txBody/a:bodyPr', 'top',    Inches(0.05), 'p:txBody/a:bodyPr'),
        ('p:txBody/a:bodyPr', 'right',  Inches(0.1),  'p:txBody/a:bodyPr'),
        ('p:txBody/a:bodyPr', 'bottom', Inches(0.05), 'p:txBody/a:bodyPr'),
    ])
    def margin_set_fixture(self, request):
        txBody_cxml, side, new_value, expected_txBody_cxml = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        prop_name = "margin_%s" % side
        expected_xml = xml(expected_txBody_cxml)
        return text_frame, prop_name, new_value, expected_xml

    @pytest.fixture(params=[
        ('p:txBody/a:p/a:r/a:t"foobar"',                     'foobar'),
        ('p:txBody/(a:p/a:r/a:t"foo",a:p/a:r/a:t"bar")',     'foo\nbar'),
        ('p:txBody/(a:p,a:p/a:r/a:t"foo",a:p/a:r/a:t"bar")', '\nfoo\nbar'),
    ])
    def text_get_fixture(self, request):
        txBody_cxml, expected_value = request.param
        text_frame = TextFrame(element(txBody_cxml), None)
        return text_frame, expected_value

    # fixture components -----------------------------------

    @pytest.fixture
    def text_frame(self, txBody):
        return TextFrame(txBody, None)

    @pytest.fixture
    def text_frame_with_parent_(self, request):
        parent_ = loose_mock(request, name='parent_')
        text_frame = TextFrame(None, parent_)
        return text_frame, parent_

    @pytest.fixture
    def txBody(self, txBody_bldr):
        return txBody_bldr.element

    @pytest.fixture
    def txBody_bldr(self):
        p_bldr = a_p()
        bodyPr_bldr = a_bodyPr()
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
        )

    @pytest.fixture
    def txBody_xml(self, txBody_bldr):
        return txBody_bldr.xml()

    @pytest.fixture
    def txBody_with_2_paras(self, _txBody_with_2_paras_bldr):
        return _txBody_with_2_paras_bldr.element

    @pytest.fixture
    def txBody_with_2_paras_xml(self, _txBody_with_2_paras_bldr):
        return _txBody_with_2_paras_bldr.xml()

    @pytest.fixture
    def txBody_with_anchor_ctr_xml(self):
        p_bldr = a_p()
        bodyPr_bldr = a_bodyPr().with_anchor('ctr')
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
                      .xml()
        )

    @pytest.fixture
    def txBody_with_text_xml(self):
        t_bldr = a_t().with_text('foobar')
        r_bldr = an_r().with_child(t_bldr)
        p_bldr = a_p().with_child(r_bldr)
        bodyPr_bldr = a_bodyPr()
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
                      .xml()
        )

    @pytest.fixture
    def txBody_with_wrap_off_xml(self):
        p_bldr = a_p()
        bodyPr_bldr = a_bodyPr().with_wrap('none')
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
                      .xml()
        )

    @pytest.fixture
    def txBody_with_wrap_on_xml(self):
        p_bldr = a_p()
        bodyPr_bldr = a_bodyPr().with_wrap('square')
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
                      .xml()
        )

    @pytest.fixture
    def _txBody_with_2_paras_bldr(self):
        p_bldr = a_p()
        bodyPr_bldr = a_bodyPr()
        return (
            a_txBody().with_nsdecls()
                      .with_child(bodyPr_bldr)
                      .with_child(p_bldr)
                      .with_child(p_bldr)
        )


class DescribeFont(object):

    def it_knows_its_bold_setting(self, bold_get_fixture):
        font, expected_value = bold_get_fixture
        assert font.bold == expected_value

    def it_can_change_its_bold_setting(self, bold_set_fixture):
        font, new_value, expected_xml = bold_set_fixture
        font.bold = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_italic_setting(self, italic_get_fixture):
        font, expected_value = italic_get_fixture
        assert font.italic == expected_value

    def it_can_change_its_italic_setting(self, italic_set_fixture):
        font, new_value, expected_xml = italic_set_fixture
        font.italic = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_underline_setting(self, underline_get_fixture):
        font, expected_value = underline_get_fixture
        assert font.underline is expected_value, 'got %s' % font.underline

    def it_can_change_its_underline_setting(self, underline_set_fixture):
        font, new_value, expected_xml = underline_set_fixture
        font.underline = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_size(self, size_get_fixture):
        font, expected_value = size_get_fixture
        assert font.size == expected_value

    def it_can_change_its_size(self, size_set_fixture):
        font, new_value, expected_xml = size_set_fixture
        font.size = new_value
        assert font._element.xml == expected_xml

    def it_knows_its_latin_typeface(self, name_get_fixture):
        font, expected_value = name_get_fixture
        assert font.name == expected_value

    def it_can_change_its_latin_typeface(self, name_set_fixture):
        font, new_value, expected_xml = name_set_fixture
        font.name = new_value
        assert font._element.xml == expected_xml

    def it_provides_access_to_its_color(self, font):
        assert isinstance(font.color, ColorFormat)

    def it_provides_access_to_its_fill(self, font):
        assert isinstance(font.fill, FillFormat)

    # fixtures ---------------------------------------------

    @pytest.fixture(params=[
        ('a:rPr',      None),
        ('a:rPr{b=0}', False),
        ('a:rPr{b=1}', True),
    ])
    def bold_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(params=[
        ('a:rPr',      True,  'a:rPr{b=1}'),
        ('a:rPr{b=1}', False, 'a:rPr{b=0}'),
        ('a:rPr{b=0}', None,  'a:rPr'),
    ])
    def bold_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[
        ('a:rPr',      None),
        ('a:rPr{i=0}', False),
        ('a:rPr{i=1}', True),
    ])
    def italic_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(params=[
        ('a:rPr',      True,  'a:rPr{i=1}'),
        ('a:rPr{i=1}', False, 'a:rPr{i=0}'),
        ('a:rPr{i=0}', None,  'a:rPr'),
    ])
    def italic_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[
        ('a:rPr',                          None),
        ('a:rPr/a:latin{typeface=Foobar}', 'Foobar'),
    ])
    def name_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(params=[
        ('a:rPr',                           'Foobar',
         'a:rPr/a:latin{typeface=Foobar}'),
        ('a:rPr/a:latin{typeface=Foobar}',  'Barfoo',
         'a:rPr/a:latin{typeface=Barfoo}'),
        ('a:rPr/a:latin{typeface=Barfoo}',  None,
         'a:rPr'),
    ])
    def name_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[
        ('a:rPr',          None),
        ('a:rPr{sz=2400}', 304800),
    ])
    def size_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(params=[
        ('a:rPr',          Pt(24), 'a:rPr{sz=2400}'),
        ('a:rPr{sz=2400}', None,   'a:rPr'),
    ])
    def size_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    @pytest.fixture(params=[
        ('a:rPr',         None),
        ('a:rPr{u=none}', False),
        ('a:rPr{u=sng}',  True),
        ('a:rPr{u=dbl}',  MSO_UNDERLINE.DOUBLE_LINE),
    ])
    def underline_get_fixture(self, request):
        rPr_cxml, expected_value = request.param
        font = Font(element(rPr_cxml))
        return font, expected_value

    @pytest.fixture(params=[
        ('a:rPr',         True,                    'a:rPr{u=sng}'),
        ('a:rPr{u=sng}',  False,                   'a:rPr{u=none}'),
        ('a:rPr{u=none}', MSO_UNDERLINE.WAVY_LINE, 'a:rPr{u=wavy}'),
        ('a:rPr{u=wavy}', MSO_UNDERLINE.NONE,      'a:rPr{u=none}'),
        ('a:rPr{u=wavy}', None,                    'a:rPr'),
    ])
    def underline_set_fixture(self, request):
        rPr_cxml, new_value, expected_rPr_cxml = request.param
        font = Font(element(rPr_cxml))
        expected_xml = xml(expected_rPr_cxml)
        return font, new_value, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def font(self):
        return Font(element('a:rPr'))


class Describe_Hyperlink(object):

    def it_knows_the_target_url_of_the_hyperlink(self, hlink_with_url_):
        hlink, rId, url = hlink_with_url_
        assert hlink.address == url
        hlink.part.target_ref.assert_called_once_with(rId)

    def it_has_None_for_address_when_no_hyperlink_is_present(self, hlink):
        assert hlink.address is None

    def it_can_set_the_target_url(
            self, hlink, rPr_with_hlinkClick_xml, url):
        hlink.address = url
        # verify -----------------------
        hlink.part.relate_to.assert_called_once_with(
            url, RT.HYPERLINK, is_external=True
        )
        assert hlink._rPr.xml == rPr_with_hlinkClick_xml
        assert hlink.address == url

    def it_can_remove_the_hyperlink(self, remove_hlink_fixture_):
        hlink, rPr_xml, rId = remove_hlink_fixture_
        hlink.address = None
        assert hlink._rPr.xml == rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId)

    def it_should_remove_the_hyperlink_when_url_set_to_empty_string(
            self, remove_hlink_fixture_):
        hlink, rPr_xml, rId = remove_hlink_fixture_
        hlink.address = ''
        assert hlink._rPr.xml == rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId)

    def it_can_change_the_target_url(self, change_hlink_fixture_):
        # fixture ----------------------
        hlink, rId_existing, new_url, new_rPr_xml = change_hlink_fixture_
        # exercise ---------------------
        hlink.address = new_url
        # verify -----------------------
        assert hlink._rPr.xml == new_rPr_xml
        hlink.part.drop_rel.assert_called_once_with(rId_existing)
        hlink.part.relate_to.assert_called_once_with(
            new_url, RT.HYPERLINK, is_external=True
        )

    # fixtures ---------------------------------------------

    @pytest.fixture
    def change_hlink_fixture_(
            self, request, hlink_with_hlinkClick, rId, rId_2, part_, url_2):
        hlinkClick_bldr = an_hlinkClick().with_rId(rId_2)
        new_rPr_xml = (
            an_rPr().with_nsdecls('a', 'r')
                    .with_child(hlinkClick_bldr)
                    .xml()
        )
        part_.relate_to.return_value = rId_2
        property_mock(request, _Hyperlink, 'part', return_value=part_)
        return hlink_with_hlinkClick, rId, url_2, new_rPr_xml

    @pytest.fixture
    def hlink(self, request, part_):
        rPr = an_rPr().with_nsdecls('a', 'r').element
        hlink = _Hyperlink(rPr, None)
        property_mock(request, _Hyperlink, 'part', return_value=part_)
        return hlink

    @pytest.fixture
    def hlink_with_hlinkClick(self, request, rPr_with_hlinkClick_bldr):
        rPr = rPr_with_hlinkClick_bldr.element
        return _Hyperlink(rPr, None)

    @pytest.fixture
    def hlink_with_url_(
            self, request, part_, hlink_with_hlinkClick, rId, url):
        property_mock(request, _Hyperlink, 'part', return_value=part_)
        return hlink_with_hlinkClick, rId, url

    @pytest.fixture
    def part_(self, request, url, rId):
        """
        Mock Part instance suitable for patching into _Hyperlink.part
        property. It returns url for target_ref() and rId for relate_to().
        """
        part_ = instance_mock(request, Part)
        part_.target_ref.return_value = url
        part_.relate_to.return_value = rId
        return part_

    @pytest.fixture
    def rId(self):
        return 'rId2'

    @pytest.fixture
    def rId_2(self):
        return 'rId6'

    @pytest.fixture
    def remove_hlink_fixture_(
            self, request, hlink_with_hlinkClick, rPr_xml, rId):
        property_mock(request, _Hyperlink, 'part')
        return hlink_with_hlinkClick, rPr_xml, rId

    @pytest.fixture
    def rPr_with_hlinkClick_bldr(self, rId):
        hlinkClick_bldr = an_hlinkClick().with_rId(rId)
        rPr_bldr = (
            an_rPr().with_nsdecls('a', 'r')
                    .with_child(hlinkClick_bldr)
        )
        return rPr_bldr

    @pytest.fixture
    def rPr_with_hlinkClick_xml(self, rPr_with_hlinkClick_bldr):
        return rPr_with_hlinkClick_bldr.xml()

    @pytest.fixture
    def rPr_xml(self):
        return an_rPr().with_nsdecls('a', 'r').xml()

    @pytest.fixture
    def url(self):
        return 'https://github.com/scanny/python-pptx'

    @pytest.fixture
    def url_2(self):
        return 'https://pypi.python.org/pypi/python-pptx'


class Describe_Paragraph(object):

    def it_knows_its_horizontal_alignment(self, alignment_get_fixture):
        paragraph, expected_value = alignment_get_fixture
        assert paragraph.alignment == expected_value

    def it_can_change_its_horizontal_alignment(self, alignment_set_fixture):
        paragraph, new_value, expected_xml = alignment_set_fixture
        paragraph.alignment = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_its_indentation_level(self, level_get_fixture):
        paragraph, expected_value = level_get_fixture
        assert paragraph.level == expected_value

    def it_can_change_its_indentation_level(self, level_set_fixture):
        paragraph, new_value, expected_xml = level_set_fixture
        paragraph.level = new_value
        assert paragraph._element.xml == expected_xml

    def it_knows_what_text_it_contains(self, text_get_fixture):
        paragraph, expected_value = text_get_fixture
        text = paragraph.text
        assert text == expected_value
        assert isinstance(text, unicode)

    def it_can_change_its_text(self, text_set_fixture):
        paragraph, new_value, expected_xml = text_set_fixture
        paragraph.text = new_value
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_its_runs(self, runs_fixture):
        paragraph, expected_text = runs_fixture
        runs = paragraph.runs
        assert tuple(r.text for r in runs) == expected_text
        for r in runs:
            assert isinstance(r, _Run)
            assert r._parent == paragraph

    def it_can_clear_itself_of_content(self, clear_fixture):
        paragraph, expected_xml = clear_fixture
        paragraph.clear()
        assert paragraph._element.xml == expected_xml

    def it_provides_access_to_the_default_paragraph_font(
            self, paragraph, Font_):
        font = paragraph.font
        Font_.assert_called_once_with(paragraph._defRPr)
        assert font == Font_.return_value

    def it_can_add_a_run(self, paragraph, p_with_r_xml):
        run = paragraph.add_run()
        assert paragraph._p.xml == p_with_r_xml
        assert isinstance(run, _Run)

    # fixtures ---------------------------------------------

    @pytest.fixture(params=[
        ('a:p',                 None),
        ('a:p/a:pPr{algn=ctr}', PP_ALIGN.CENTER),
        ('a:p/a:pPr{algn=r}',   PP_ALIGN.RIGHT),
    ])
    def alignment_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(params=[
        ('a:p',                  PP_ALIGN.LEFT,    'a:p/a:pPr{algn=l}'),
        ('a:p/a:pPr{algn=l}',    PP_ALIGN.JUSTIFY, 'a:p/a:pPr{algn=just}'),
        ('a:p/a:pPr{algn=just}', None,             'a:p/a:pPr'),
    ])
    def alignment_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture(params=[
        ('a:p/a:r/a:t"foo"',              'a:p'),
        ('a:p/(a:br,a:r/a:t"foo")',       'a:p'),
        ('a:p/(a:fld,a:br,a:r/a:t"foo")', 'a:p'),
    ])
    def clear_fixture(self, request):
        p_cxml, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, expected_xml

    @pytest.fixture(params=[
        ('a:p',              0),
        ('a:p/a:pPr{lvl=2}', 2),
    ])
    def level_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(params=[
        ('a:p',              1, 'a:p/a:pPr{lvl=1}'),
        ('a:p/a:pPr{lvl=1}', 2, 'a:p/a:pPr{lvl=2}'),
        ('a:p/a:pPr{lvl=2}', 0, 'a:p/a:pPr'),
    ])
    def level_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    @pytest.fixture
    def runs_fixture(self):
        p_cxml = 'a:p/(a:r/a:t"Foo",a:r/a:t"Bar",a:r/a:t"Baz")'
        paragraph = _Paragraph(element(p_cxml), None)
        expected_text = ('Foo', 'Bar', 'Baz')
        return paragraph, expected_text

    @pytest.fixture(params=[
        ('a:p/a:r/a:t"foobar"',                               'foobar'),
        ('a:p/(a:r/a:t"foo", a:r/a:t"bar")',                  'foobar'),
        ('a:p/(a:r/a:t"foo", a:br, a:r/a:t"bar")',            'foo\nbar'),
        ('a:p/(a:r/a:t"foo ", a:fld/a:t"42", a:r/a:t" bar")', 'foo 42 bar'),
        ('a:p/(a:r/a:t" foo", a:br, a:fld/a:t"42")',          ' foo\n42'),
        ('a:p/(a:pPr,a:r/a:t"foobar",a:endParaRPr)',          'foobar'),
        ('a:p/a:fld/a:t"42"',                                 '42'),
        ('a:p/a:br',                                          '\n'),
    ])
    def text_get_fixture(self, request):
        p_cxml, expected_value = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        return paragraph, expected_value

    @pytest.fixture(params=[
        ('a:p/(a:r/a:t"foo",a:r/a:t"bar")', 'foobar', 'a:p/a:r/a:t"foobar"'),
        ('a:p', 'foo\nbar',      'a:p/(a:r/a:t"foo",a:br,a:r/a:t"bar")'),
        ('a:p', '\nfoo\n',       'a:p/(a:br,a:r/a:t"foo",a:br)'),
        ('a:p', 'foo\n',         'a:p/(a:r/a:t"foo",a:br)'),
        ('a:p', '7-bit str',     'a:p/a:r/a:t"7-bit str"'),
        ('a:p', '8-ɓïȶ str',    u'a:p/a:r/a:t"8-ɓïȶ str"'),
        ('a:p', u'ŮŦƑ literal', u'a:p/a:r/a:t"ŮŦƑ literal"'),
        ('a:p', unicode('utf-8 unicode: Hér er texti', 'utf-8'),
         u'a:p/a:r/a:t"utf-8 unicode: Hér er texti"'),
    ])
    def text_set_fixture(self, request):
        p_cxml, new_value, expected_p_cxml = request.param
        paragraph = _Paragraph(element(p_cxml), None)
        expected_xml = xml(expected_p_cxml)
        return paragraph, new_value, expected_xml

    # fixture components -----------------------------------

    @pytest.fixture
    def Font_(self, request):
        return class_mock(request, 'pptx.text.Font')

    @pytest.fixture
    def p_bldr(self):
        return a_p().with_nsdecls()

    @pytest.fixture
    def p_with_r_xml(self):
        run_bldr = an_r().with_child(a_t())
        return a_p().with_nsdecls().with_child(run_bldr).xml()

    @pytest.fixture
    def paragraph(self, p_bldr):
        return _Paragraph(p_bldr.element, None)


class Describe_Run(object):

    def it_provides_access_to_its_font(self, font_fixture):
        run, rPr, Font_, font_ = font_fixture
        font = run.font
        Font_.assert_called_once_with(rPr)
        assert font == font_

    def it_provides_access_to_a_hyperlink_proxy(self, hyperlink_fixture):
        run, rPr, _Hyperlink_, hlink_ = hyperlink_fixture
        hlink = run.hyperlink
        _Hyperlink_.assert_called_once_with(rPr, run)
        assert hlink is hlink_

    def it_can_get_the_text_of_the_run(self, text_get_fixture):
        run, expected_value = text_get_fixture
        text = run.text
        assert text == expected_value
        assert isinstance(text, unicode)

    def it_can_change_its_text(self, text_set_fixture):
        run, new_value, expected_xml = text_set_fixture
        run.text = new_value
        assert run._r.xml == expected_xml

    # fixtures ---------------------------------------------

    @pytest.fixture
    def font_fixture(self, Font_, font_):
        r = element('a:r/a:rPr')
        rPr = r.rPr
        run = _Run(r, None)
        return run, rPr, Font_, font_

    @pytest.fixture
    def hyperlink_fixture(self, _Hyperlink_, hlink_):
        r = element('a:r/a:rPr')
        rPr = r.rPr
        run = _Run(r, None)
        return run, rPr, _Hyperlink_, hlink_

    @pytest.fixture
    def text_get_fixture(self):
        r = element('a:r/a:t"foobar"')
        run = _Run(r, None)
        return run, 'foobar'

    @pytest.fixture(params=[
        ('a:r/a:t', 'barfoo', 'a:r/a:t"barfoo"'),
    ])
    def text_set_fixture(self, request):
        r_cxml, new_value, expected_r_cxml = request.param
        run = _Run(element(r_cxml), None)
        expected_xml = xml(expected_r_cxml)
        return run, new_value, expected_xml

    # fixture components -----------------------------------

    @pytest.fixture
    def Font_(self, request, font_):
        return class_mock(request, 'pptx.text.Font', return_value=font_)

    @pytest.fixture
    def font_(self, request):
        return instance_mock(request, Font)

    @pytest.fixture
    def _Hyperlink_(self, request, hlink_):
        return class_mock(
            request, 'pptx.text._Hyperlink', return_value=hlink_
        )

    @pytest.fixture
    def hlink_(self, request):
        return instance_mock(request, _Hyperlink)
