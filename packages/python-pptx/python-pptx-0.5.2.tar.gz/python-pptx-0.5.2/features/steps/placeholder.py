# encoding: utf-8

"""
Gherkin step implementations for placeholder-related features.
"""

from __future__ import absolute_import

from behave import given, when, then

from pptx import Presentation

from .helpers import saved_pptx_path, test_pptx, test_text


# given ===================================================

@given('a bullet body placeholder')
def given_a_bullet_body_placeholder(context):
    context.prs = Presentation()
    slide_layout = context.prs.slide_layouts[1]
    context.sld = context.prs.slides.add_slide(slide_layout)
    context.body = context.sld.shapes.placeholders[1]


@given('a layout placeholder having directly set position and size')
def given_layout_placeholder_with_directly_set_pos_and_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slide_layouts[0].placeholders[1]


@given('a layout placeholder having no direct position or size settings')
def given_layout_placeholder_with_no_direct_pos_or_size_settings(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slide_layouts[0].placeholders[0]


@given('a master placeholder')
def given_a_master_placeholder(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.master_placeholder = prs.slide_master.placeholders[1]


@given('a slide placeholder having directly set position and size')
def given_slide_placeholder_with_directly_set_pos_and_size(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[0].placeholders[1]


@given('a slide placeholder having no direct position or size settings')
def given_slide_placeholder_with_no_direct_pos_or_size_settings(context):
    prs = Presentation(test_pptx('ph-inherit-props'))
    context.placeholder = prs.slides[0].placeholders[0]


# when ====================================================

@when('I indent the first paragraph')
def step_when_indent_first_paragraph(context):
    context.body.text_frame.paragraphs[0].level = 1


@when("I set the title text of the slide")
def step_when_set_slide_title_text(context):
    context.slide.shapes.title.text = test_text


# then ====================================================

@then('I can get the placeholder dimensions')
def then_I_can_get_the_placeholder_dimensions(context):
    placeholder = context.master_placeholder
    assert placeholder.width == 6923112, 'got %d' % placeholder.width
    assert placeholder.height == 3484984, 'got %d' % placeholder.height


@then('I can get the placeholder position')
def then_I_can_get_the_placeholder_position(context):
    placeholder = context.master_placeholder
    assert placeholder.left == 1110444, 'got %d' % placeholder.left
    assert placeholder.top == 1686508, 'got %d' % placeholder.top


@then('I get the direct settings when I query position and size')
def then_I_get_direct_settings_when_query_pos_and_size(context):
    placeholder = context.placeholder
    assert placeholder.left == 468312, 'got %s' % placeholder.left
    assert placeholder.top == 1700212, 'got %s' % placeholder.top
    assert placeholder.width == 8208143, 'got %s' % placeholder.width
    assert placeholder.height == 4537099, 'got %s' % placeholder.height


@then('I get inherited settings when I query position and size')
def then_I_get_inherited_settings_when_I_query_position_and_size(context):
    placeholder = context.placeholder
    assert placeholder.left == 457200, 'got %s' % placeholder.left
    assert placeholder.top == 274638, 'got %s' % placeholder.top
    assert placeholder.width == 8229600, 'got %s' % placeholder.width
    assert placeholder.height == 1143000, 'got %s' % placeholder.height


@then('the paragraph is indented')
def then_paragraph_is_indented(context):
    prs = Presentation(saved_pptx_path)
    sld = prs.slides[0]
    body = sld.shapes.placeholders[1]
    p = body.text_frame.paragraphs[0]
    assert p.level == 1


@then('the text appears in the title placeholder')
def step_then_text_appears_in_title_placeholder(context):
    prs = Presentation(saved_pptx_path)
    title_shape = prs.slides[0].shapes.title
    title_text = title_shape.text_frame.paragraphs[0].runs[0].text
    assert title_text == test_text
